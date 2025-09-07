
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typing import List, Tuple

def extract_slide_text_and_images(pptx_path: str):
    prs = Presentation(pptx_path)
    all_slides = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        image_blobs: List[Tuple[str, bytes]] = []  # (name, blob)
        for shp in slide.shapes:
            if hasattr(shp, "has_text_frame") and shp.has_text_frame:
                tx = "\n".join([p.text for p in shp.text_frame.paragraphs])
                if tx:
                    texts.append(tx)
            if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shp.image
                    ext = image.ext or "bin"
                    name = f"slide{idx}_{image.filename or 'image'}.{ext}"
                    image_blobs.append((name, image.blob))
                except Exception:
                    pass
        text_content = "\n".join(texts).strip()
        all_slides.append((idx, text_content, image_blobs))
    return all_slides
